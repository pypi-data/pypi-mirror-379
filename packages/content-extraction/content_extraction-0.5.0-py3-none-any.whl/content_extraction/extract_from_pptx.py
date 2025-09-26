import os
import argparse
import sys
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


logger = logging.getLogger(__name__)


def extract_content(pptx_path: str, output_dir: str) -> str | None:
    logger.info(f'Extracting content from {pptx_path} to {output_dir}')
    """
    Extracts text, tables, and images from a PPTX file and saves them
    into an HTML file and an images directory.

    Args:
        pptx_path (str): The path to the input PowerPoint presentation.
        output_dir (str): The directory where the output HTML and images will be saved.

    Returns:
        tuple[str, str] | tuple[None, None]: A tuple containing the path to the
        output HTML file and the images directory, or (None, None) on failure.
    """
    logger.info(f'[Extracting content] started from {pptx_path}')
    images_dir = os.path.join(output_dir, 'images')
    html_out_path = os.path.join(output_dir, 'index.html')
    logger.debug(f'[Extracting content] Images directory: {images_dir}, HTML output path: {html_out_path}')

    # Ensure output directories exist
    os.makedirs(images_dir, exist_ok=True)

    try:
        prs = Presentation(pptx_path)
    except Exception:
        logger.error(f'[Extracting content] Could not open or parse {pptx_path}')
        return None

    html_lines = [
        '<!DOCTYPE html>',
        '<html lang="en">',
        '<head>',
        '  <meta charset="UTF-8">',
        '  <title>Extracted PPTX Content</title>',
        '</head>',
        '<body>',
    ]

    image_counter = 0

    logger.debug('[Extracting content] stated processing slides')
    for slide_idx, slide in enumerate(prs.slides, start=1):
        logger.debug(f'[Extracting content] stated processing slide {slide_idx}')
        html_lines.append(f'  <section id="slide-{slide_idx}">')

        # 1) Title (if any)
        title_text = None
        if slide.has_notes_slide:  # Check for title in shapes first
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    title_text = shape.text_frame.text.strip()
                    break
        if title_text:
            html_lines.append(f'    <h1>{title_text}</h1>')

        # 2) Walk every shape
        for shape in slide.shapes:
            # -- TABLES --
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                html_lines.append('    <table border="1">')
                table = shape.table
                for row in table.rows:
                    html_lines.append('      <tr>')
                    for cell in row.cells:
                        cell_txt = cell.text.replace('\n', '<br/>')
                        html_lines.append(f'        <td>{cell_txt}</td>')
                    html_lines.append('      </tr>')
                html_lines.append('    </table>')

            # -- IMAGES --
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_counter += 1
                ext = image.ext  # e.g. 'png', 'jpeg'
                img_name = f'slide{slide_idx}_img{image_counter}.{ext}'
                img_path = os.path.join(images_dir, img_name)
                with open(img_path, 'wb') as f:
                    f.write(image.blob)
                # Relative path for the src attribute
                html_lines.append(f'    <img src="images/{img_name}" alt="Slide {slide_idx} image"/>')

            # -- TEXT (including bullets) --
            elif shape.has_text_frame:
                # skip re-printing the title placeholder
                if shape.is_placeholder and shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.SUBTITLE,
                ):
                    continue

                in_list = False
                for para in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in para.runs).strip()
                    if not text:
                        continue

                    # any indent > 0 treat as a bullet
                    if para.level > 0:
                        if not in_list:
                            html_lines.append('    <ul>')
                            in_list = True
                        html_lines.append(f'      <li>{text}</li>')
                    else:
                        if in_list:
                            html_lines.append('    </ul>')
                            in_list = False
                        html_lines.append(f'    <p>{text}</p>')

                if in_list:
                    html_lines.append('    </ul>')

        html_lines.append('  </section>')
        logger.debug(f'[Extracting content] completed processing slide {slide_idx}')

    logger.debug('[Extracting content] completed processing slides')
    html_lines.extend(['</body>', '</html>'])

    # Write out the final HTML file
    logger.debug('[Extracting content] started saving HTML file')
    with open(html_out_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(html_lines))
    logger.debug('[Extracting content] completed saving HTML file')

    logger.info(f'[Extracting content] completed from {pptx_path}')
    return html_out_path


def main():
    """Main function to handle command line arguments and execute the script."""
    parser = argparse.ArgumentParser(
        description='Extract content from a PowerPoint (PPTX) file to HTML.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.pptx                      # Outputs to 'output/' directory by default
  %(prog)s presentation.pptx -o extracted_content # Outputs to 'extracted_content/' directory
""",
    )
    parser.add_argument('pptx_file', help='Path to the input PPTX file.')
    parser.add_argument(
        '-o',
        '--output',
        default='output',
        help="Path to the output directory (if not provided, defaults to 'output').",
    )
    args = parser.parse_args()

    if not os.path.exists(args.pptx_file):
        logger.error(f'Input file not found at {args.pptx_file}')
        return 1

    html_out, images_out = extract_content(args.pptx_file, args.output)

    if html_out and images_out:
        logger.info(f'Successfully extracted content to {html_out} with images in {images_out}/')
        return 0

    logger.error(f'Extraction failed for {args.pptx_file}.')
    return 1


if __name__ == '__main__':
    sys.exit(main())
